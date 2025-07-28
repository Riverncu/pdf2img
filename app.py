import streamlit as st
import fitz
from io import BytesIO
from PIL import Image
import zipfile
from pptx import Presentation
from pptx.util import Inches

VALID_CODES = st.secrets["employee_codes"]

def convert_pdf_to_pptx(doc):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for i in range(doc.page_count):
        page = doc.load_page(i)
        pix = page.get_pixmap(dpi=150)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        img_path = f"/tmp/page_{i + 1}.png"
        img.save(img_path)

        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

def render_preview_and_download(doc, pdf_name, start_page, end_page, output_format, dpi, conversion_format):
    if conversion_format == "PPTX":
        pptx_file = convert_pdf_to_pptx(doc)
        st.download_button("⬇️ Download PPTX", pptx_file, file_name="converted.pptx",
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        return

    progress_bar = st.progress(0)
    zip_buffer = BytesIO()
    with zipfile.ZipFile(zip_buffer, "w") as zip_file:
        container = st.container()
        page_indices = range(start_page - 1, end_page)
        cols_per_row = 4
        rows = (len(page_indices) + cols_per_row - 1) // cols_per_row

        for row_idx in range(rows):
            cols = container.columns(cols_per_row)
            for col_idx in range(cols_per_row):
                page_idx = row_idx * cols_per_row + col_idx
                if page_idx >= len(page_indices):
                    break

                i = page_indices[page_idx]
                page = doc.load_page(i)
                pix = page.get_pixmap(dpi=dpi)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

                img_buffer = BytesIO()
                img.save(img_buffer, format=output_format)
                img_bytes = img_buffer.getvalue()

                with cols[col_idx]:
                    st.image(img, caption=f"Page {i + 1}", use_container_width=True)
                    st.download_button(
                        label=f"⬇️ Download Page {i + 1} as {output_format}",
                        data=img_bytes,
                        file_name=f"page_{i + 1}.{output_format.lower()}",
                        mime=f"image/{output_format.lower()}",
                    )

                zip_file.writestr(f"page_{i + 1}.{output_format.lower()}", img_bytes)

                progress_bar.progress((page_idx + 1) / len(page_indices))

    progress_bar.empty()
    zip_buffer.seek(0)

    st.download_button(
        label=f"📦 Download ZIP of pages {start_page} to {end_page}",
        data=zip_buffer,
        file_name=f"{pdf_name}_pages_{start_page}_to_{end_page}.zip",
        mime="application/zip",
    )


# --- Custom CSS ---
st.markdown("""
<style>
.header {
    font-size: 2.8rem;
    font-weight: 800;
    background: linear-gradient(90deg, #0072ff, #00c6ff);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    margin-bottom: 30px;
    text-align: center;
}
.sidebar .sidebar-content {
    background-color: #e6f7ff;
    padding: 25px;
    border-radius: 15px;
    box-shadow: 0 6px 15px rgb(0 198 255 / 0.3);
    border: 1px solid #00aaff;
}
.footer {
    font-size: 0.95rem;
    color: #005073;
    text-align: center;
    padding: 25px;
    margin-top: 60px;
    border-top: 2px solid #00aaff;
    background: #e0f7ff;
    border-radius: 10px 10px 0 0;
}
div.stDownloadButton > button {
    background-color: #0072ff;
    color: white;
    font-weight: 700;
    border-radius: 12px;
    padding: 10px 20px;
    border: none;
    box-shadow: 0 4px 8px rgb(0 114 255 / 0.6);
    transition: background-color 0.3s ease, box-shadow 0.3s ease;
}
div.stDownloadButton > button:hover {
    background-color: #0057cc;
    box-shadow: 0 6px 14px rgb(0 87 204 / 0.8);
}
.stImage > img {
    border-radius: 10px;
    box-shadow: 0 5px 15px rgb(0 114 255 / 0.3);
    margin-bottom: 8px;
}
</style>
""", unsafe_allow_html=True)

st.set_page_config(page_title="PDF to Image Converter", layout="wide")

st.markdown('<div class="header">📄 PDF to Image Converter</div>', unsafe_allow_html=True)

if "authorized" not in st.session_state:
    st.session_state.authorized = False
    st.session_state.user_greeting = ""

if not st.session_state.authorized:
    code_input = st.text_input("Please enter last 3 digits of your employee ID to access", max_chars=3)

    if code_input in VALID_CODES:
        st.session_state.authorized = True
        st.session_state.user_greeting = VALID_CODES[code_input]
        st.rerun()
    elif code_input != "":
        st.warning("⚠️ You are not authorized to use this app. Please enter valid employee code.")

else:
    st.success(st.session_state.user_greeting)

    with st.sidebar:
        st.markdown('<div class="sidebar">', unsafe_allow_html=True)
        pdf_file = st.file_uploader("Upload PDF file", type=["pdf"])
        if pdf_file is not None:
            doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
            total_pages = doc.page_count

            start_page = st.number_input("Start page", min_value=1, max_value=total_pages, value=1)
            end_page = st.number_input("End page", min_value=1, max_value=total_pages, value=total_pages)
            output_format = st.selectbox("Select output image format", ["PNG", "JPEG"])
            quality = st.radio("Select image quality (DPI)", ["Normal (300)", "HighQ (600)"])
            conversion_format = st.selectbox("Convert PDF to other format", ["None", "PPTX"])
        else:
            start_page = end_page = output_format = quality = conversion_format = None

        st.markdown("</div>", unsafe_allow_html=True)

    # Ở main area hiện preview và download khi đã có file và các lựa chọn hợp lệ
    if pdf_file is not None and start_page is not None and end_page is not None and output_format is not None and quality is not None and conversion_format is not None:
        if start_page <= end_page:
            dpi = 300 if "300" in quality else 600
            render_preview_and_download(doc, pdf_file.name, start_page, end_page, output_format, dpi, conversion_format)
        else:
            st.error("❌ Start page must be less or equal to End page")

    # Footer
    st.markdown("""
    <div class="footer">
        Developed by FTC-Harper  — © 2025
    </div>
    """, unsafe_allow_html=True)
